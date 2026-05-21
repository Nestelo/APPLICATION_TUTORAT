"""
Générateurs pour les rapports en différents formats
Support: CSV, Excel, Word, PDF, PowerPoint
"""

import csv
import json
from datetime import datetime
from io import StringIO, BytesIO

from django.http import HttpResponse
from django.template.loader import render_to_string

# Essayer d'importer les dépendances optionnelles
try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
    from reportlab.lib import colors
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches, Pt as PPTPt
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class RapportExporter:
    """Classe de base pour l'export de rapports"""
    
    def __init__(self, rapport_type, period, data):
        self.rapport_type = rapport_type
        self.period = period
        self.data = data
        self.titre = f'Rapport {rapport_type.title()}'
        self.date_generation = datetime.now().strftime('%d/%m/%Y à %H:%M:%S')
    
    def export_csv(self):
        """Exporte le rapport en CSV"""
        output = StringIO()
        writer = csv.writer(output)
        
        # En-tête
        writer.writerow([self.titre])
        writer.writerow([f'Période: {self.period}'])
        writer.writerow([f'Généré le: {self.date_generation}'])
        writer.writerow([])
        
        # Contenu selon le type
        if self.rapport_type == 'utilisateurs':
            self._export_csv_utilisateurs(writer)
        elif self.rapport_type == 'tutorat':
            self._export_csv_tutorat(writer)
        elif self.rapport_type == 'ressources':
            self._export_csv_ressources(writer)
        elif self.rapport_type == 'forum':
            self._export_csv_forum(writer)
        
        # Préparer la réponse
        output.seek(0)
        response = HttpResponse(output.getvalue(), content_type='text/csv; charset=utf-8')
        response['Content-Disposition'] = f'attachment; filename=rapport_{self.rapport_type}_{self.period}.csv'
        return response
    
    def _export_csv_utilisateurs(self, writer):
        """Export CSV spécifique pour les utilisateurs"""
        writer.writerow(['STATISTIQUES GLOBALES'])
        writer.writerow(['Total utilisateurs', self.data.get('total_users', 0)])
        writer.writerow(['Taux de croissance (%)', self.data.get('growth_rate', 0)])
        writer.writerow(['Taux d\'activité (%)', self.data.get('activity_rate', 0)])
        writer.writerow([])
        
        # Répartition par rôle
        writer.writerow(['RÉPARTITION PAR RÔLE'])
        writer.writerow(['Rôle', 'Nombre'])
        for item in self.data.get('users_by_role', []):
            writer.writerow([item.get('role', 'N/A'), item.get('count', 0)])
        writer.writerow([])
        
        # Liste détaillée par rôle
        writer.writerow(['DÉTAIL PAR RÔLE'])
        writer.writerow(['Email', 'Nom', 'Prénom', 'Rôle', 'Filière', 'Année', 'Date inscription'])
        for user in self.data.get('recent_users', []):
            writer.writerow([
                user.get('email', 'N/A'),
                user.get('nom', 'N/A'),
                user.get('prenom', 'N/A'),
                user.get('role', 'N/A'),
                user.get('filiere', 'N/A'),
                user.get('annee', 'N/A'),
                user.get('date_inscription', 'N/A'),
            ])
        writer.writerow([])
        
        # Distribution par filière
        writer.writerow(['DISTRIBUTION PAR FILIÈRE'])
        writer.writerow(['Filière', 'Nombre'])
        for item in self.data.get('users_by_filiere', []):
            writer.writerow([item.get('filiere', 'N/A'), item.get('count', 0)])
    
    def _export_csv_tutorat(self, writer):
        """Export CSV spécifique pour le tutorat"""
        writer.writerow(['STATISTIQUES GLOBALES'])
        writer.writerow(['Total séances', self.data.get('total_seances', 0)])
        writer.writerow(['Taux de completion (%)', self.data.get('completion_rate', 0)])
        writer.writerow(['Satisfaction moyenne', self.data.get('average_satisfaction', 0)])
        writer.writerow([])
        
        # Matières populaires
        writer.writerow(['MATIÈRES POPULAIRES'])
        writer.writerow(['Matière', 'Nombre séances', 'Tarif moyen'])
        for item in self.data.get('matieres_populaires', []):
            writer.writerow([
                item.get('matiere', 'N/A'),
                item.get('count', 0),
                item.get('avg_tarif', 0),
            ])
        writer.writerow([])
        
        # Tuteurs actifs
        writer.writerow(['TUTEURS LES PLUS ACTIFS'])
        writer.writerow(['Nom', 'Prénom', 'Email', 'Nombre séances'])
        for item in self.data.get('tuteurs_actifs', []):
            writer.writerow([
                item.get('nom', 'N/A'),
                item.get('prenom', 'N/A'),
                item.get('email', 'N/A'),
                item.get('seance_count', 0),
            ])
    
    def _export_csv_ressources(self, writer):
        """Export CSV spécifique pour les ressources"""
        writer.writerow(['STATISTIQUES GLOBALES'])
        writer.writerow(['Total ressources', self.data.get('total_ressources', 0)])
        writer.writerow(['Ressources publiées', self.data.get('published_ressources', 0)])
        writer.writerow(['Taux de publication (%)', self.data.get('publication_rate', 0)])
        writer.writerow(['Total téléchargements', self.data.get('total_downloads', 0)])
        writer.writerow(['Total vues', self.data.get('total_views', 0)])
        writer.writerow([])
        
        # Ressources les plus téléchargées
        writer.writerow(['RESSOURCES LES PLUS TÉLÉCHARGÉES'])
        writer.writerow(['Titre', 'Matière', 'Téléchargements', 'Vues'])
        for item in self.data.get('most_downloaded', []):
            writer.writerow([
                item.get('titre', 'N/A'),
                item.get('matiere', 'N/A'),
                item.get('nb_telechargements', 0),
                item.get('nb_vues', 0),
            ])
        writer.writerow([])
        
        # Auteurs productifs
        writer.writerow(['AUTEURS LES PLUS PRODUCTIFS'])
        writer.writerow(['Nom', 'Prénom', 'Email', 'Nombre ressources'])
        for item in self.data.get('prolific_authors', []):
            writer.writerow([
                item.get('nom', 'N/A'),
                item.get('prenom', 'N/A'),
                item.get('email', 'N/A'),
                item.get('resource_count', 0),
            ])
    
    def _export_csv_forum(self, writer):
        """Export CSV spécifique pour le forum"""
        writer.writerow(['STATISTIQUES GLOBALES'])
        writer.writerow(['Total questions', self.data.get('total_questions', 0)])
        writer.writerow(['Questions résolues', self.data.get('resolved_questions', 0)])
        writer.writerow(['Taux de résolution (%)', self.data.get('resolution_rate', 0)])
        writer.writerow([])
        
        # Questions populaires
        writer.writerow(['QUESTIONS LES PLUS VUES'])
        writer.writerow(['Titre', 'Vues', 'Réponses'])
        for item in self.data.get('most_viewed', []):
            writer.writerow([
                item.get('titre', 'N/A'),
                item.get('nb_vues', 0),
                item.get('nb_reponses', 0),
            ])
    
    def export_excel(self):
        """Exporte le rapport en Excel"""
        if not HAS_OPENPYXL:
            return HttpResponse(
                'Openpyxl non installé. Installez avec: pip install openpyxl',
                status=500
            )
        
        wb = Workbook()
        ws = wb.active
        ws.title = 'Rapport'
        
        # Styles
        header_font = Font(bold=True, color="FFFFFF", size=12)
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        subheader_font = Font(bold=True, size=11)
        subheader_fill = PatternFill(start_color="B4C7E7", end_color="B4C7E7", fill_type="solid")
        
        # En-tête
        row = 1
        ws[f'A{row}'] = self.titre
        ws[f'A{row}'].font = Font(bold=True, size=14)
        row += 1
        ws[f'A{row}'] = f'Période: {self.period}'
        row += 1
        ws[f'A{row}'] = f'Généré le: {self.date_generation}'
        row += 2
        
        # Contenu selon le type
        if self.rapport_type == 'utilisateurs':
            row = self._export_excel_utilisateurs(ws, row, header_font, header_fill, subheader_font, subheader_fill)
        elif self.rapport_type == 'tutorat':
            row = self._export_excel_tutorat(ws, row, header_font, header_fill, subheader_font, subheader_fill)
        elif self.rapport_type == 'ressources':
            row = self._export_excel_ressources(ws, row, header_font, header_fill, subheader_font, subheader_fill)
        elif self.rapport_type == 'forum':
            row = self._export_excel_forum(ws, row, header_font, header_fill, subheader_font, subheader_fill)
        
        # Ajuster les largeurs de colonnes
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # Préparer la réponse
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        
        response = HttpResponse(
            output,
            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
        response['Content-Disposition'] = f'attachment; filename=rapport_{self.rapport_type}_{self.period}.xlsx'
        return response
    
    def _export_excel_utilisateurs(self, ws, start_row, header_font, header_fill, subheader_font, subheader_fill):
        """Export Excel spécifique pour les utilisateurs"""
        row = start_row
        
        # Statistiques globales
        ws[f'A{row}'] = 'STATISTIQUES GLOBALES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        ws[f'A{row}'] = 'Total utilisateurs'
        ws[f'B{row}'] = self.data.get('total_users', 0)
        row += 1
        
        ws[f'A{row}'] = 'Taux de croissance (%)'
        ws[f'B{row}'] = self.data.get('growth_rate', 0)
        row += 1
        
        ws[f'A{row}'] = 'Taux d\'activité (%)'
        ws[f'B{row}'] = self.data.get('activity_rate', 0)
        row += 2
        
        # Répartition par rôle
        ws[f'A{row}'] = 'RÉPARTITION PAR RÔLE'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        ws[f'A{row}'] = 'Rôle'
        ws[f'B{row}'] = 'Nombre'
        ws[f'A{row}'].font = header_font
        ws[f'B{row}'].font = header_font
        ws[f'A{row}'].fill = header_fill
        ws[f'B{row}'].fill = header_fill
        row += 1
        
        for item in self.data.get('users_by_role', []):
            ws[f'A{row}'] = item.get('role', 'N/A')
            ws[f'B{row}'] = item.get('count', 0)
            row += 1
        
        row += 1
        
        # Détail par rôle
        ws[f'A{row}'] = 'DÉTAIL DES UTILISATEURS'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        headers = ['Email', 'Nom', 'Prénom', 'Rôle', 'Filière', 'Année', 'Date inscription']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
        row += 1
        
        for user in self.data.get('recent_users', []):
            ws[f'A{row}'] = user.get('email', 'N/A')
            ws[f'B{row}'] = user.get('nom', 'N/A')
            ws[f'C{row}'] = user.get('prenom', 'N/A')
            ws[f'D{row}'] = user.get('role', 'N/A')
            ws[f'E{row}'] = user.get('filiere', 'N/A')
            ws[f'F{row}'] = user.get('annee', 'N/A')
            ws[f'G{row}'] = user.get('date_inscription', 'N/A')
            row += 1
        
        return row
    
    def _export_excel_tutorat(self, ws, start_row, header_font, header_fill, subheader_font, subheader_fill):
        """Export Excel spécifique pour le tutorat"""
        row = start_row
        
        # Statistiques globales
        ws[f'A{row}'] = 'STATISTIQUES GLOBALES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        ws[f'A{row}'] = 'Total séances'
        ws[f'B{row}'] = self.data.get('total_seances', 0)
        row += 1
        
        ws[f'A{row}'] = 'Taux de completion (%)'
        ws[f'B{row}'] = self.data.get('completion_rate', 0)
        row += 1
        
        ws[f'A{row}'] = 'Satisfaction moyenne'
        ws[f'B{row}'] = self.data.get('average_satisfaction', 0)
        row += 2
        
        # Matières populaires
        ws[f'A{row}'] = 'MATIÈRES POPULAIRES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        headers = ['Matière', 'Nombre séances', 'Tarif moyen']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
        row += 1
        
        for item in self.data.get('matieres_populaires', []):
            ws[f'A{row}'] = item.get('matiere', 'N/A')
            ws[f'B{row}'] = item.get('count', 0)
            ws[f'C{row}'] = item.get('avg_tarif', 0)
            row += 1
        
        row += 1
        
        # Tuteurs actifs
        ws[f'A{row}'] = 'TUTEURS LES PLUS ACTIFS'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        headers = ['Nom', 'Prénom', 'Email', 'Nombre séances']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
        row += 1
        
        for item in self.data.get('tuteurs_actifs', []):
            ws[f'A{row}'] = item.get('nom', 'N/A')
            ws[f'B{row}'] = item.get('prenom', 'N/A')
            ws[f'C{row}'] = item.get('email', 'N/A')
            ws[f'D{row}'] = item.get('seance_count', 0)
            row += 1
        
        return row
    
    def _export_excel_ressources(self, ws, start_row, header_font, header_fill, subheader_font, subheader_fill):
        """Export Excel spécifique pour les ressources"""
        row = start_row
        
        # Statistiques globales
        ws[f'A{row}'] = 'STATISTIQUES GLOBALES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        ws[f'A{row}'] = 'Total ressources'
        ws[f'B{row}'] = self.data.get('total_ressources', 0)
        row += 1
        
        ws[f'A{row}'] = 'Ressources publiées'
        ws[f'B{row}'] = self.data.get('published_ressources', 0)
        row += 1
        
        ws[f'A{row}'] = 'Total téléchargements'
        ws[f'B{row}'] = self.data.get('total_downloads', 0)
        row += 1
        
        ws[f'A{row}'] = 'Total vues'
        ws[f'B{row}'] = self.data.get('total_views', 0)
        row += 2
        
        # Ressources les plus téléchargées
        ws[f'A{row}'] = 'RESSOURCES LES PLUS TÉLÉCHARGÉES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        headers = ['Titre', 'Matière', 'Téléchargements', 'Vues']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
        row += 1
        
        for item in self.data.get('most_downloaded', []):
            ws[f'A{row}'] = item.get('titre', 'N/A')
            ws[f'B{row}'] = item.get('matiere', 'N/A')
            ws[f'C{row}'] = item.get('nb_telechargements', 0)
            ws[f'D{row}'] = item.get('nb_vues', 0)
            row += 1
        
        return row
    
    def _export_excel_forum(self, ws, start_row, header_font, header_fill, subheader_font, subheader_fill):
        """Export Excel spécifique pour le forum"""
        row = start_row
        
        # Statistiques globales
        ws[f'A{row}'] = 'STATISTIQUES GLOBALES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        ws[f'A{row}'] = 'Total questions'
        ws[f'B{row}'] = self.data.get('total_questions', 0)
        row += 1
        
        ws[f'A{row}'] = 'Questions résolues'
        ws[f'B{row}'] = self.data.get('resolved_questions', 0)
        row += 1
        
        ws[f'A{row}'] = 'Taux de résolution (%)'
        ws[f'B{row}'] = self.data.get('resolution_rate', 0)
        row += 2
        
        # Questions les plus vues
        ws[f'A{row}'] = 'QUESTIONS LES PLUS VUES'
        ws[f'A{row}'].font = subheader_font
        ws[f'A{row}'].fill = subheader_fill
        row += 1
        
        headers = ['Titre', 'Vues', 'Réponses']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
        row += 1
        
        for item in self.data.get('most_viewed', []):
            ws[f'A{row}'] = item.get('titre', 'N/A')
            ws[f'B{row}'] = item.get('nb_vues', 0)
            ws[f'C{row}'] = item.get('nb_reponses', 0)
            row += 1
        
        return row
    
    def export_word(self):
        """Exporte le rapport en Word"""
        if not HAS_DOCX:
            return HttpResponse(
                'python-docx non installé. Installez avec: pip install python-docx',
                status=500
            )
        
        doc = Document()
        
        # En-tête
        title = doc.add_heading(self.titre, level=1)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph(f'Période: {self.period}')
        doc.add_paragraph(f'Généré le: {self.date_generation}')
        doc.add_paragraph()
        
        # Contenu selon le type
        if self.rapport_type == 'utilisateurs':
            self._export_word_utilisateurs(doc)
        elif self.rapport_type == 'tutorat':
            self._export_word_tutorat(doc)
        elif self.rapport_type == 'ressources':
            self._export_word_ressources(doc)
        elif self.rapport_type == 'forum':
            self._export_word_forum(doc)
        
        # Préparer la réponse
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        
        response = HttpResponse(
            output,
            content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        response['Content-Disposition'] = f'attachment; filename=rapport_{self.rapport_type}_{self.period}.docx'
        return response
    
    def _export_word_utilisateurs(self, doc):
        """Export Word spécifique pour les utilisateurs"""
        # Statistiques globales
        doc.add_heading('Statistiques Globales', level=2)
        table = doc.add_table(rows=4, cols=2)
        table.style = 'Light Grid Accent 1'
        
        table.cell(0, 0).text = 'Total utilisateurs'
        table.cell(0, 1).text = str(self.data.get('total_users', 0))
        table.cell(1, 0).text = 'Taux de croissance (%)'
        table.cell(1, 1).text = str(self.data.get('growth_rate', 0))
        table.cell(2, 0).text = 'Taux d\'activité (%)'
        table.cell(2, 1).text = str(self.data.get('activity_rate', 0))
        table.cell(3, 0).text = 'Période'
        table.cell(3, 1).text = self.period
        
        # Répartition par rôle
        doc.add_heading('Répartition par Rôle', level=2)
        roles_table = doc.add_table(rows=len(self.data.get('users_by_role', [])) + 1, cols=2)
        roles_table.style = 'Light Grid Accent 1'
        
        roles_table.cell(0, 0).text = 'Rôle'
        roles_table.cell(0, 1).text = 'Nombre'
        
        for idx, item in enumerate(self.data.get('users_by_role', []), 1):
            roles_table.cell(idx, 0).text = item.get('role', 'N/A')
            roles_table.cell(idx, 1).text = str(item.get('count', 0))
        
        doc.add_paragraph()
    
    def _export_word_tutorat(self, doc):
        """Export Word spécifique pour le tutorat"""
        # Statistiques globales
        doc.add_heading('Statistiques Globales', level=2)
        table = doc.add_table(rows=4, cols=2)
        table.style = 'Light Grid Accent 1'
        
        table.cell(0, 0).text = 'Total séances'
        table.cell(0, 1).text = str(self.data.get('total_seances', 0))
        table.cell(1, 0).text = 'Taux de completion (%)'
        table.cell(1, 1).text = str(self.data.get('completion_rate', 0))
        table.cell(2, 0).text = 'Satisfaction moyenne'
        table.cell(2, 1).text = str(self.data.get('average_satisfaction', 0))
        table.cell(3, 0).text = 'Période'
        table.cell(3, 1).text = self.period
        
        doc.add_paragraph()
    
    def _export_word_ressources(self, doc):
        """Export Word spécifique pour les ressources"""
        # Statistiques globales
        doc.add_heading('Statistiques Globales', level=2)
        table = doc.add_table(rows=5, cols=2)
        table.style = 'Light Grid Accent 1'
        
        table.cell(0, 0).text = 'Total ressources'
        table.cell(0, 1).text = str(self.data.get('total_ressources', 0))
        table.cell(1, 0).text = 'Ressources publiées'
        table.cell(1, 1).text = str(self.data.get('published_ressources', 0))
        table.cell(2, 0).text = 'Total téléchargements'
        table.cell(2, 1).text = str(self.data.get('total_downloads', 0))
        table.cell(3, 0).text = 'Total vues'
        table.cell(3, 1).text = str(self.data.get('total_views', 0))
        table.cell(4, 0).text = 'Période'
        table.cell(4, 1).text = self.period
        
        doc.add_paragraph()
    
    def _export_word_forum(self, doc):
        """Export Word spécifique pour le forum"""
        # Statistiques globales
        doc.add_heading('Statistiques Globales', level=2)
        table = doc.add_table(rows=4, cols=2)
        table.style = 'Light Grid Accent 1'
        
        table.cell(0, 0).text = 'Total questions'
        table.cell(0, 1).text = str(self.data.get('total_questions', 0))
        table.cell(1, 0).text = 'Questions résolues'
        table.cell(1, 1).text = str(self.data.get('resolved_questions', 0))
        table.cell(2, 0).text = 'Taux de résolution (%)'
        table.cell(2, 1).text = str(self.data.get('resolution_rate', 0))
        table.cell(3, 0).text = 'Période'
        table.cell(3, 1).text = self.period
        
        doc.add_paragraph()
    
    def export_pdf(self):
        """Exporte le rapport en PDF"""
        if not HAS_REPORTLAB:
            return HttpResponse(
                'reportlab non installé. Installez avec: pip install reportlab',
                status=500
            )
        
        output = BytesIO()
        doc = SimpleDocTemplate(output, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # En-tête
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=colors.HexColor('#1F4E78'),
            spaceAfter=12,
            alignment=1  # Center
        )
        
        story.append(Paragraph(self.titre, title_style))
        story.append(Paragraph(f'<i>Période: {self.period}</i>', styles['Normal']))
        story.append(Paragraph(f'<i>Généré le: {self.date_generation}</i>', styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))
        
        # Contenu selon le type
        if self.rapport_type == 'utilisateurs':
            self._export_pdf_utilisateurs(story, styles)
        elif self.rapport_type == 'tutorat':
            self._export_pdf_tutorat(story, styles)
        elif self.rapport_type == 'ressources':
            self._export_pdf_ressources(story, styles)
        elif self.rapport_type == 'forum':
            self._export_pdf_forum(story, styles)
        
        doc.build(story)
        output.seek(0)
        
        response = HttpResponse(output, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename=rapport_{self.rapport_type}_{self.period}.pdf'
        return response
    
    def _export_pdf_utilisateurs(self, story, styles):
        """Export PDF spécifique pour les utilisateurs"""
        # Statistiques globales
        story.append(Paragraph('<b>Statistiques Globales</b>', styles['Heading2']))
        data = [
            ['Total utilisateurs', str(self.data.get('total_users', 0))],
            ['Taux de croissance (%)', str(self.data.get('growth_rate', 0))],
            ['Taux d\'activité (%)', str(self.data.get('activity_rate', 0))],
        ]
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.3 * inch))
    
    def _export_pdf_tutorat(self, story, styles):
        """Export PDF spécifique pour le tutorat"""
        # Statistiques globales
        story.append(Paragraph('<b>Statistiques Globales</b>', styles['Heading2']))
        data = [
            ['Total séances', str(self.data.get('total_seances', 0))],
            ['Taux de completion (%)', str(self.data.get('completion_rate', 0))],
            ['Satisfaction moyenne', str(self.data.get('average_satisfaction', 0))],
        ]
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.3 * inch))
    
    def _export_pdf_ressources(self, story, styles):
        """Export PDF spécifique pour les ressources"""
        # Statistiques globales
        story.append(Paragraph('<b>Statistiques Globales</b>', styles['Heading2']))
        data = [
            ['Total ressources', str(self.data.get('total_ressources', 0))],
            ['Ressources publiées', str(self.data.get('published_ressources', 0))],
            ['Total téléchargements', str(self.data.get('total_downloads', 0))],
            ['Total vues', str(self.data.get('total_views', 0))],
        ]
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.3 * inch))
    
    def _export_pdf_forum(self, story, styles):
        """Export PDF spécifique pour le forum"""
        # Statistiques globales
        story.append(Paragraph('<b>Statistiques Globales</b>', styles['Heading2']))
        data = [
            ['Total questions', str(self.data.get('total_questions', 0))],
            ['Questions résolues', str(self.data.get('resolved_questions', 0))],
            ['Taux de résolution (%)', str(self.data.get('resolution_rate', 0))],
        ]
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(t)
        story.append(Spacer(1, 0.3 * inch))
    
    def export_powerpoint(self):
        """Exporte le rapport en PowerPoint"""
        if not HAS_PPTX:
            return HttpResponse(
                'python-pptx non installé. Installez avec: pip install python-pptx',
                status=500
            )
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Slide 1: Titre
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.add_textbox(PPTInches(0.5), PPTInches(2.5), PPTInches(9), PPTInches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = self.titre
        p.font.size = PPTPt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(PPTInches(0.5), PPTInches(5), PPTInches(9), PPTInches(2))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = f"Période: {self.period}\nGénéré le: {self.date_generation}"
        p.font.size = PPTPt(20)
        p.alignment = PP_ALIGN.CENTER
        
        # Contenu selon le type
        if self.rapport_type == 'utilisateurs':
            self._export_pptx_utilisateurs(prs)
        elif self.rapport_type == 'tutorat':
            self._export_pptx_tutorat(prs)
        elif self.rapport_type == 'ressources':
            self._export_pptx_ressources(prs)
        elif self.rapport_type == 'forum':
            self._export_pptx_forum(prs)
        
        # Préparer la réponse
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        response = HttpResponse(
            output,
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = f'attachment; filename=rapport_{self.rapport_type}_{self.period}.pptx'
        return response
    
    def _export_pptx_utilisateurs(self, prs):
        """Export PowerPoint spécifique pour les utilisateurs"""
        # Slide: Statistiques
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Statistiques Globales"
        
        content = slide.placeholders[1].text_frame
        content.text = f"Total utilisateurs: {self.data.get('total_users', 0)}\n"
        content.text += f"Taux de croissance: {self.data.get('growth_rate', 0)}%\n"
        content.text += f"Taux d'activité: {self.data.get('activity_rate', 0)}%"
    
    def _export_pptx_tutorat(self, prs):
        """Export PowerPoint spécifique pour le tutorat"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Statistiques Globales"
        
        content = slide.placeholders[1].text_frame
        content.text = f"Total séances: {self.data.get('total_seances', 0)}\n"
        content.text += f"Taux de completion: {self.data.get('completion_rate', 0)}%\n"
        content.text += f"Satisfaction moyenne: {self.data.get('average_satisfaction', 0)}/5"
    
    def _export_pptx_ressources(self, prs):
        """Export PowerPoint spécifique pour les ressources"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Statistiques Globales"
        
        content = slide.placeholders[1].text_frame
        content.text = f"Total ressources: {self.data.get('total_ressources', 0)}\n"
        content.text += f"Ressources publiées: {self.data.get('published_ressources', 0)}\n"
        content.text += f"Total téléchargements: {self.data.get('total_downloads', 0)}\n"
        content.text += f"Total vues: {self.data.get('total_views', 0)}"
    
    def _export_pptx_forum(self, prs):
        """Export PowerPoint spécifique pour le forum"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Statistiques Globales"
        
        content = slide.placeholders[1].text_frame
        content.text = f"Total questions: {self.data.get('total_questions', 0)}\n"
        content.text += f"Questions résolues: {self.data.get('resolved_questions', 0)}\n"
        content.text += f"Taux de résolution: {self.data.get('resolution_rate', 0)}%"
